import argparse
import json
import os
from datetime import datetime

from pptx import Presentation


def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--ref", required=True)
    return parser.parse_args()


def all_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texts.append(para.text)
    return " ".join(texts)


def main():
    args = parse_args()
    errors = []
    warnings = []

    # 파일 존재 & 크기
    if not os.path.exists(args.pptx):
        errors.append(f"파일 없음: {args.pptx}")
        write_report(args, 0, 0, 0, False, errors, warnings)
        return

    file_size_kb = os.path.getsize(args.pptx) // 1024
    if file_size_kb < 10:
        errors.append(f"파일 크기 미달: {file_size_kb}KB < 10KB")

    # ref JSON 로드
    with open(args.ref, "r", encoding="utf-8") as f:
        ref_data = json.load(f)

    # PPT 열기
    try:
        prs = Presentation(args.pptx)
    except Exception as e:
        errors.append(f"pptx 열기 실패: {e}")
        write_report(args, file_size_kb, 0, 0, False, errors, warnings)
        return

    total_slides = len(prs.slides)
    types = set(r.get("type", "other") for r in ref_data)
    expected = 1 + 1 + len(types) + len(ref_data) + 1  # 표지+목차+섹션+개별+출처

    if abs(total_slides - expected) > 1:
        errors.append(f"슬라이드 수 불일치: 실제={total_slides}, 예상={expected}")

    # 표지 슬라이드 텍스트 확인
    if total_slides > 0:
        cover_text = all_text(prs.slides[0])
        if "건설업 사고사례 보고서" not in cover_text:
            errors.append("표지 슬라이드에 '건설업 사고사례 보고서' 텍스트 없음")

    # 출처 슬라이드 확인
    if total_slides > 0:
        src_text = all_text(prs.slides[-1])
        if "KOSHA" not in src_text and "한국산업안전보건공단" not in src_text:
            warnings.append("출처 슬라이드에 KOSHA 텍스트 없음")

    # 빈 슬라이드 검사
    for i, slide in enumerate(prs.slides):
        text = all_text(slide).strip()
        if not text:
            warnings.append(f"슬라이드 {i+1}: 텍스트 없음 (빈 슬라이드)")

    passed = len(errors) == 0
    write_report(args, file_size_kb, total_slides, expected, passed, errors, warnings)


def write_report(args, file_size_kb, total_slides, expected_slides, passed, errors, warnings):
    date_str = os.path.basename(args.pptx).replace(".pptx", "")
    output_dir = os.path.join("output", "qa")
    os.makedirs(output_dir, exist_ok=True)
    report_path = os.path.join(output_dir, f"{date_str}_report.json")

    report = {
        "pptx_path": args.pptx,
        "file_size_kb": file_size_kb,
        "total_slides": total_slides,
        "expected_slides": expected_slides,
        "passed": passed,
        "errors": errors,
        "warnings": warnings,
        "checked_at": datetime.now().isoformat(),
    }
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)

    status = "PASS" if passed else "FAIL"
    print(f"[QA] {status} — errors={len(errors)}, warnings={len(warnings)}: {report_path}")
    if not passed:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
